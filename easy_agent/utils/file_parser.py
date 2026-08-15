"""File content parser - extracts text from various file formats"""

import csv
import logging
import os
import subprocess
from pathlib import Path

try:
    import fitz
except ImportError:
    fitz = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except ImportError:
    pdfminer_extract_text = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import olefile
except ImportError:
    olefile = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

logger = logging.getLogger(__name__)

MAX_FILE_SIZE = 50 * 1024 * 1024

TEXT_EXTENSIONS = {
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".cs",
    ".rb",
    ".php",
    ".swift",
    ".kt",
    ".kts",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".sql",
    ".r",
    ".m",
    ".yaml",
    ".yml",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".css",
    ".scss",
    ".less",
    ".vue",
    ".svelte",
    ".md",
    ".rst",
    ".txt",
    ".ini",
    ".cfg",
    ".conf",
    ".toml",
    ".env",
    ".gitignore",
    ".dockerfile",
    ".proto",
    ".graphql",
    ".tf",
    ".gradle",
    ".properties",
    ".log",
    ".csv",
    ".tsv",
    ".tex",
    ".bib",
    ".makefile",
    ".cmake",
    ".nim",
    ".zig",
    ".v",
    ".vhdl",
    ".sv",
    ".lua",
    ".pl",
    ".pm",
    ".tcl",
    ".dart",
    ".elm",
    ".erl",
    ".hrl",
    ".ex",
    ".exs",
    ".clj",
    ".cljs",
    ".edn",
    ".hs",
    ".lhs",
    ".fs",
    ".fsx",
    ".ml",
    ".mli",
    ".jl",
    ".rkt",
    ".scm",
    ".ss",
    ".coffee",
    ".litcoffee",
    ".styl",
    ".sass",
    ".pug",
    ".jade",
    ".haml",
    ".slim",
    ".twig",
    ".blade",
    ".erb",
    ".ejs",
    ".mustache",
    ".handlebars",
    ".njk",
}


def parse_file_content(file_path: str, mime_type: str = "") -> str:
    ext = os.path.splitext(file_path)[1].lower()
    path = Path(file_path)

    if not path.exists():
        logger.warning(f"文件不存在: {file_path}")
        return ""

    if path.stat().st_size > MAX_FILE_SIZE:
        logger.warning(f"文件过大，跳过解析: {file_path} ({path.stat().st_size} bytes)")
        return f"[文件过大无法解析: {path.name} ({_format_size(path.stat().st_size)})]"

    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        elif ext == ".docx":
            return _parse_docx(path)
        elif ext == ".doc":
            return _parse_doc(path)
        elif ext == ".xlsx":
            return _parse_xlsx(path)
        elif ext == ".xls":
            return _parse_xls(path)
        elif ext == ".pptx":
            return _parse_pptx(path)
        elif ext == ".ppt":
            return _parse_ppt(path)
        elif ext == ".rtf":
            return _parse_rtf(path)
        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"):
            return _parse_image(path)
        elif ext == ".csv":
            return _parse_csv(path)
        elif ext in TEXT_EXTENSIONS:
            return _parse_text(path)
        else:
            return _parse_text(path)
    except Exception as e:
        logger.warning(f"解析文件失败 {file_path}: {e}")
        return f"[文件解析失败: {e}]"


def _format_size(size: int) -> str:
    if size < 1024:
        return f"{size} B"
    elif size < 1024 * 1024:
        return f"{size / 1024:.1f} KB"
    else:
        return f"{size / (1024 * 1024):.1f} MB"


def _parse_text(path: Path) -> str:
    encodings = [
        "utf-8",
        "gbk",
        "gb2312",
        "latin-1",
        "shift-jis",
        "big5",
        "euc-kr",
        "cp1252",
    ]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="latin-1", errors="replace")


def _parse_pdf(path: Path) -> str:
    try:
        if fitz is None:
            raise ImportError

        doc = fitz.open(str(path))
        text_parts = []
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text)
        doc.close()
        text = "\n".join(text_parts)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"PyMuPDF 解析失败: {e}")

    try:
        if pdfminer_extract_text is None:
            raise ImportError

        text = pdfminer_extract_text(str(path))
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"pdfminer 解析失败: {e}")

    try:
        if pdfplumber is None:
            raise ImportError

        with pdfplumber.open(str(path)) as pdf:
            text_parts = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            text = "\n".join(text_parts)
            if text.strip():
                return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"pdfplumber 解析失败: {e}")

    return f"[PDF文件: {path.name}，未安装PDF解析库，请安装: uv pip install pymupdf pdfminer.six pdfplumber]"


def _parse_docx(path: Path) -> str:
    try:
        if Document is None:
            raise ImportError

        doc = Document(str(path))
        text_parts = []

        for p in doc.paragraphs:
            if p.text.strip():
                text_parts.append(p.text)

        for table in doc.tables:
            text_parts.append("--- [表格] ---")
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        text = "\n".join(text_parts)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"python-docx 解析失败: {e}")

    return (
        f"[DOCX文件: {path.name}，未安装python-docx，请安装: uv pip install python-docx]"
    )


def _parse_doc(path: Path) -> str:
    try:
        result = subprocess.run(
            ["antiword", str(path)], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning(f"antiword 解析失败: {e}")

    try:
        result = subprocess.run(
            ["catdoc", str(path)], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning(f"catdoc 解析失败: {e}")

    try:
        if olefile is None:
            raise ImportError

        with olefile.OleFileIO(str(path)) as ole:
            if ole.exists("WordDocument"):
                stream = ole.openstream("1Table")
                data = stream.read()
                text = data.decode("utf-8", errors="replace")
                if text.strip():
                    return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"olefile 解析失败: {e}")

    return f"[DOC文件(旧格式): {path.name}，建议转换为docx格式，或安装: apt install antiword catdoc]"


def _parse_xlsx(path: Path) -> str:
    try:
        if openpyxl is None:
            raise ImportError

        wb = openpyxl.load_workbook(str(path), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                if row_text.strip():
                    text_parts.append(row_text)
        text = "\n".join(text_parts)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"openpyxl 解析失败: {e}")

    return f"[XLSX文件: {path.name}，未安装openpyxl，请安装: uv pip install openpyxl]"


def _parse_xls(path: Path) -> str:
    try:
        if xlrd is None:
            raise ImportError

        wb = xlrd.open_workbook(str(path))
        text_parts = []
        for sheet in wb.sheets():
            text_parts.append(f"--- Sheet: {sheet.name} ---")
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_text = "\t".join(str(v) if v else "" for v in row_values)
                if row_text.strip():
                    text_parts.append(row_text)
        text = "\n".join(text_parts)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"xlrd 解析失败: {e}")

    return (
        f"[XLS文件(旧格式): {path.name}，建议转换为xlsx格式，或安装: uv pip install xlrd]"
    )


def _parse_pptx(path: Path) -> str:
    try:
        if Presentation is None:
            raise ImportError

        prs = Presentation(str(path))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            if slide_texts:
                text_parts.append(f"--- Slide {i + 1} ---")
                text_parts.extend(slide_texts)
        text = "\n".join(text_parts)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"python-pptx 解析失败: {e}")

    return (
        f"[PPTX文件: {path.name}，未安装python-pptx，请安装: uv pip install python-pptx]"
    )


def _parse_ppt(path: Path) -> str:
    try:
        result = subprocess.run(
            ["catppt", str(path)], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning(f"catppt 解析失败: {e}")

    return f"[PPT文件(旧格式): {path.name}，建议转换为pptx格式，或安装: apt install catdoc]"


def _parse_rtf(path: Path) -> str:
    try:
        if rtf_to_text is None:
            raise ImportError

        with open(path, "r", encoding="utf-8", errors="replace") as f:
            rtf_text = f.read()
        text = rtf_to_text(rtf_text)
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"striprtf 解析失败: {e}")

    try:
        result = subprocess.run(
            ["unrtf", "--text", str(path)], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning(f"unrtf 解析失败: {e}")

    return _parse_text(path)


def _parse_image(path: Path) -> str:
    try:
        if Image is None or pytesseract is None:
            raise ImportError

        img = Image.open(str(path))
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"OCR 解析失败: {e}")

    return (
        f"[图片文件: {path.name}，未安装OCR库，请安装: uv pip install pytesseract pillow]"
    )


def _parse_csv(path: Path) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
            if rows:
                return "\n".join(",".join(row) for row in rows)
    except UnicodeDecodeError:
        pass
    except Exception as e:
        logger.warning(f"CSV 解析失败: {e}")

    try:
        with open(path, "r", encoding="gbk") as f:
            reader = csv.reader(f)
            rows = list(reader)
            if rows:
                return "\n".join(",".join(row) for row in rows)
    except Exception:
        pass

    return _parse_text(path)
