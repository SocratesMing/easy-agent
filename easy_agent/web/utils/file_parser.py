"""File content parser - extracts text from various file formats"""

import logging
import os
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def parse_file_content(file_path: str, mime_type: str = "") -> str:
    """Extract text content from a file based on its type.

    Supports: PDF, DOCX, XLSX, images (OCR via pytesseract), plain text, code, CSV.
    """
    ext = os.path.splitext(file_path)[1].lower()
    path = Path(file_path)

    if not path.exists():
        logger.warning(f"文件不存在: {file_path}")
        return ""

    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        elif ext in (".docx", ".doc"):
            return _parse_docx(path)
        elif ext in (".xlsx", ".xls"):
            return _parse_xlsx(path)
        elif ext in (".pptx", ".ppt"):
            return _parse_pptx(path)
        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
            return _parse_image(path)
        elif ext in (".csv",):
            return _parse_csv(path)
        elif ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs",
                     ".c", ".cpp", ".h", ".hpp", ".cs", ".rb", ".php", ".swift",
                     ".kt", ".scala", ".sh", ".bash", ".zsh", ".sql", ".r", ".m",
                     ".yaml", ".yml", ".json", ".xml", ".html", ".css", ".scss",
                     ".less", ".vue", ".svelte", ".md", ".rst", ".txt", ".ini",
                     ".cfg", ".conf", ".toml", ".env", ".gitignore", ".dockerfile",
                     ".proto", ".graphql", ".tf", ".gradle", ".properties"):
            return _parse_text(path)
        else:
            # 尝试以文本方式读取未知格式
            return _parse_text(path)
    except Exception as e:
        logger.warning(f"解析文件失败 {file_path}: {e}")
        return f"[文件解析失败: {e}]"


def _parse_text(path: Path) -> str:
    """Read plain text file with encoding auto-detection."""
    encodings = ["utf-8", "gbk", "gb2312", "latin-1", "shift-jis", "big5"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="latin-1", errors="replace")


def _parse_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF or pdfminer."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"PyMuPDF 解析失败: {e}")

    try:
        from pdfminer.high_level import extract_text
        text = extract_text(str(path))
        if text.strip():
            return text
    except ImportError:
        pass

    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            if text.strip():
                return text
    except ImportError:
        pass

    return f"[PDF文件: {path.name}，未安装PDF解析库，请安装: pip install pymupdf pdfminer.six pdfplumber]"


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs)
    except ImportError:
        return f"[DOCX文件: {path.name}，请安装: pip install python-docx]"


def _parse_xlsx(path: Path) -> str:
    """Extract text from XLSX."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_str.strip():
                    lines.append(row_str)
        wb.close()
        return "\n".join(lines)
    except ImportError:
        pass

    try:
        import pandas as pd
        dfs = pd.read_excel(str(path), sheet_name=None)
        lines = []
        for sheet_name, df in dfs.items():
            lines.append(f"=== Sheet: {sheet_name} ===")
            lines.append(df.to_string(index=False))
        return "\n".join(lines)
    except ImportError:
        return f"[XLSX文件: {path.name}，请安装: pip install openpyxl pandas]"


def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)
    except ImportError:
        return f"[PPTX文件: {path.name}，请安装: pip install python-pptx]"


def _parse_image(path: Path) -> str:
    """Extract text from image via OCR (pytesseract)."""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if text.strip():
            return text
    except ImportError:
        pass

    try:
        import easyocr
        reader = easyocr.Reader(["ch_sim", "en"], gpu=False)
        result = reader.readtext(str(path), detail=0)
        return "\n".join(result)
    except ImportError:
        pass

    return f"[图片文件: {path.name}，未安装OCR库，请安装: pip install pillow pytesseract easyocr]"


def _parse_csv(path: Path) -> str:
    """Extract text from CSV."""
    try:
        import csv
        import io
        with open(str(path), encoding="utf-8", errors="replace") as f:
            content = f.read()

        lines = []
        reader = csv.reader(io.StringIO(content))
        for row in reader:
            lines.append(" | ".join(row))
        return "\n".join(lines)
    except Exception as e:
        return _parse_text(path)
